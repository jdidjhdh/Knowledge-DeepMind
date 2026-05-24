import asyncio
import json
import os
import re
import uuid
import tempfile
import logging
from typing import Optional
from fastapi import UploadFile

import httpx
from bs4 import BeautifulSoup

from config import Settings
from models import IngestionTask, DocumentChunk, DocumentType
from .multimedia_enhancer import MultimediaEnhancer
from .code_analyzer import analyze_code_file, detect_language, CODE_SEMANTIC_PROMPT
from .speech_processor import SpeechProcessor
from .image_processor import get_image_processor

logger = logging.getLogger(__name__)


class FileHandler:
    def __init__(self, settings: Settings):
        self.settings = settings
        self.upload_dir = settings.upload_dir
        self._enhancer = None
        self._speech_processor = None
        self._image_processor = None
        os.makedirs(self.upload_dir, exist_ok=True)

    @property
    def enhancer(self):
        if self._enhancer is None:
            self._enhancer = MultimediaEnhancer(self.settings)
        return self._enhancer

    @property
    def speech_processor(self):
        if self._speech_processor is None:
            self._speech_processor = SpeechProcessor(
                settings=self.settings,
            )
        return self._speech_processor

    async def _get_image_processor(self):
        if self._image_processor is None:
            self._image_processor = await get_image_processor(settings=self.settings)
        return self._image_processor

    async def process_file(self, file: UploadFile, file_type: DocumentType) -> IngestionTask:
        task_id = str(uuid.uuid4())
        task = IngestionTask(
            task_id=task_id,
            file_path=file.filename or "",
            file_type=file_type,
            status="processing",
            progress=0.0,
        )
        try:
            content = await file.read()
            file_size = len(content)
            max_size = self.settings.max_upload_size_mb * 1024 * 1024
            if file_size > max_size:
                task.status = "failed"
                task.error = f"文件大小超过限制 ({self.settings.max_upload_size_mb}MB)"
                return task

            save_path = os.path.join(self.upload_dir, f"{task_id}_{file.filename}")
            with open(save_path, "wb") as f:
                f.write(content)

            task.progress = 0.3

            if file_type == DocumentType.CODE:
                chunks, code_analysis = await self._parse_code_with_analysis(save_path, file.filename or "")
                task.progress = 0.9
                task.status = "completed"
                task.progress = 1.0
                task.result = {
                    "chunks": [chunk.model_dump() for chunk in chunks],
                    "file_size": file_size,
                    "saved_path": save_path,
                    "code_analysis": code_analysis,
                }
            else:
                chunks = await self._extract_chunks(save_path, file.filename or "", file_type)
                task.progress = 0.9
                task.status = "completed"
                task.progress = 1.0
                task.result = {
                    "chunks": [chunk.model_dump() for chunk in chunks],
                    "file_size": file_size,
                    "saved_path": save_path,
                }
            logger.info(f"文件处理完成: {file.filename}, 提取 {len(chunks)} 个文档片段")

        except Exception as e:
            logger.error(f"文件处理失败: {file.filename}, 错误: {str(e)}")
            task.status = "failed"
            task.error = str(e)

        return task

    async def process_url(self, url: str) -> IngestionTask:
        task_id = str(uuid.uuid4())
        task = IngestionTask(
            task_id=task_id,
            file_path=url,
            file_type=DocumentType.WEB,
            status="processing",
            progress=0.0,
        )
        try:
            async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
                response = await client.get(
                    url,
                    headers={
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                                      "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
                    },
                )
                response.raise_for_status()
                html_content = response.text

            task.progress = 0.4

            chunks = self._parse_html_content(html_content, url)
            task.progress = 0.9

            save_path = os.path.join(self.upload_dir, f"{task_id}_webpage.html")
            with open(save_path, "w", encoding="utf-8") as f:
                f.write(html_content)

            task.status = "completed"
            task.progress = 1.0
            task.result = {
                "chunks": [chunk.model_dump() for chunk in chunks],
                "file_size": len(html_content),
                "saved_path": save_path,
                "url": url,
            }
            logger.info(f"URL抓取完成: {url}, 提取 {len(chunks)} 个文档片段")

        except httpx.HTTPStatusError as e:
            logger.error(f"URL请求失败: {url}, HTTP {e.response.status_code}")
            task.status = "failed"
            task.error = f"HTTP {e.response.status_code}: 网页请求失败"
        except httpx.RequestError as e:
            logger.error(f"URL连接失败: {url}, 错误: {str(e)}")
            task.status = "failed"
            task.error = f"无法连接到该网址: {str(e)}"
        except Exception as e:
            logger.error(f"URL处理失败: {url}, 错误: {str(e)}")
            task.status = "failed"
            task.error = str(e)

        return task

    def _parse_html_content(self, html_content: str, source_url: str) -> list[DocumentChunk]:
        soup = BeautifulSoup(html_content, "lxml")
        title = ""
        title_tag = soup.find("title")
        if title_tag:
            title = title_tag.get_text(strip=True)
        for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip() and len(line.strip()) > 2]

        chunks = []
        if title:
            chunks.append(DocumentChunk(
                content=f"页面标题: {title}\n\n" + "\n".join(lines),
                source_path=source_url,
                source_type=DocumentType.WEB,
                chunk_index=0,
                metadata={"title": title, "url": source_url},
            ))
        elif lines:
            chunks.append(DocumentChunk(
                content="\n".join(lines),
                source_path=source_url,
                source_type=DocumentType.WEB,
                chunk_index=0,
                metadata={"url": source_url},
            ))
        return chunks

    async def _extract_chunks(
        self, file_path: str, filename: str, file_type: DocumentType
    ) -> list[DocumentChunk]:
        if file_type == DocumentType.PDF:
            return await self._parse_pdf(file_path, filename)
        elif file_type == DocumentType.WORD:
            return await self._parse_word(file_path, filename)
        elif file_type == DocumentType.PPT:
            return await self._parse_ppt(file_path, filename)
        elif file_type == DocumentType.IMAGE:
            return await self._parse_image(file_path, filename)
        elif file_type == DocumentType.WEB:
            return await self._parse_web(file_path, filename)
        elif file_type == DocumentType.TABLE:
            return await self._parse_table(file_path, filename)
        elif file_type == DocumentType.CODE:
            return await self._parse_code(file_path, filename)
        elif file_type == DocumentType.TEXT:
            return await self._parse_text(file_path, filename)
        elif file_type == DocumentType.VIDEO:
            return await self._parse_video(file_path, filename)
        elif file_type == DocumentType.AUDIO:
            return await self._parse_audio(file_path, filename)
        else:
            return await self._parse_text(file_path, filename)

    async def _parse_pdf(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            import fitz
            chunks = []
            doc = fitz.open(file_path)
            has_ocr = False
            try:
                import pytesseract
                from PIL import Image
                has_ocr = True
            except ImportError:
                pass

            for i, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    chunks.append(DocumentChunk(
                        content=text.strip(),
                        source_path=filename,
                        source_type=DocumentType.PDF,
                        chunk_index=i,
                        metadata={"page": i + 1, "total_pages": len(doc), "source": "embedded_text"},
                    ))
                elif has_ocr:
                    pix = page.get_pixmap(dpi=200)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_text = pytesseract.image_to_string(img, lang="chi_sim+eng")
                    if ocr_text.strip():
                        chunks.append(DocumentChunk(
                            content=ocr_text.strip(),
                            source_path=filename,
                            source_type=DocumentType.PDF,
                            chunk_index=i,
                            metadata={"page": i + 1, "total_pages": len(doc), "source": "ocr"},
                        ))
            doc.close()
            return chunks
        except ImportError:
            return await self._fallback_text(file_path, filename, DocumentType.PDF)

    async def _parse_word(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            from docx import Document
            doc = Document(file_path)
            chunks = []
            full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if full_text:
                chunks.append(DocumentChunk(
                    content=full_text,
                    source_path=filename,
                    source_type=DocumentType.WORD,
                    chunk_index=0,
                    metadata={"paragraphs": len(doc.paragraphs)},
                ))
            return chunks
        except ImportError:
            return await self._fallback_text(file_path, filename, DocumentType.WORD)

    async def _parse_ppt(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            chunks = []
            for i, slide in enumerate(prs.slides):
                texts = []
                table_count = 0
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        table_text = [f"[表格 {table_count + 1}]"]
                        for row_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            table_text.append(" | ".join(cells))
                        texts.append("\n".join(table_text))
                        table_count += 1
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                texts.append(paragraph.text.strip())
                if texts:
                    chunks.append(DocumentChunk(
                        content="\n".join(texts),
                        source_path=filename,
                        source_type=DocumentType.PPT,
                        chunk_index=i,
                        metadata={
                            "slide": i + 1,
                            "total_slides": len(prs.slides),
                            "table_count": table_count,
                        },
                    ))
            return chunks
        except ImportError:
            return await self._fallback_text(file_path, filename, DocumentType.PPT)

    async def _parse_image(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            processor = await self._get_image_processor()
            result = await processor.process(file_path, filename)
        except Exception as e:
            logger.error(f"[图片增强v5] 流水线失败: {e}, 回退到基础模式")
            result = {"classification": {}, "ocr": {}, "fusion": {}, "visual_analysis": {}, "knowledge_chunks": []}

        classification = result.get("classification", {})
        ocr_result = result.get("ocr", {})
        fusion = result.get("fusion", {})
        visual_analysis = result.get("visual_analysis", {})
        knowledge_chunks = result.get("knowledge_chunks", [])

        chunks = []
        for kc in knowledge_chunks:
            chunks.append(DocumentChunk(
                content=kc["content"],
                source_path=filename,
                source_type=DocumentType.IMAGE,
                chunk_index=kc.get("chunk_index", 0),
                metadata=kc.get("metadata", {}),
            ))

        if not chunks:
            fallback_text = ocr_result.get("text", "")
            if not fallback_text:
                from PIL import Image
                img = Image.open(file_path)
                fallback_text = (
                    f"[图片] 类型: {classification.get('category', '未知')}\n"
                    f"分辨率: {img.width}x{img.height}\n"
                    f"分析: {classification.get('reason', '')}"
                )
            chunks.append(DocumentChunk(
                content=fallback_text,
                source_path=filename,
                source_type=DocumentType.IMAGE,
                chunk_index=0,
                metadata={
                    "category": classification.get("category", "photo"),
                    "ocr_engine": ocr_result.get("engine", "none"),
                    "extraction_mode": "enhanced_v5_fallback",
                },
            ))

        logger.info(
            f"[图片增强v5] 完成: {filename}, {len(chunks)} chunks, "
            f"类别={classification.get('category', '?')}, "
            f"OCR引擎={ocr_result.get('engine', '?')}, "
            f"OCR置信度={ocr_result.get('confidence', 0):.3f}, "
            f"低置信度={len(fusion.get('low_confidence_regions', []))}"
        )
        return chunks

    async def _parse_web(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            soup = BeautifulSoup(content, "lxml")
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            chunks = []
            if lines:
                chunks.append(DocumentChunk(
                    content="\n".join(lines),
                    source_path=filename,
                    source_type=DocumentType.WEB,
                    chunk_index=0,
                ))
            return chunks
        except ImportError:
            return await self._fallback_text(file_path, filename, DocumentType.WEB)

    async def _parse_table(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            import pandas as pd
            ext = os.path.splitext(file_path)[1].lower()
            if ext == ".csv":
                df = pd.read_csv(file_path)
            elif ext in (".xlsx", ".xls"):
                df = pd.read_excel(file_path)
            else:
                return await self._fallback_text(file_path, filename, DocumentType.TABLE)

            chunks = []
            description = f"表格包含 {len(df)} 行, {len(df.columns)} 列\n"
            description += f"列名: {', '.join(df.columns.astype(str))}\n"
            description += f"数据类型:\n{df.dtypes.to_string()}\n\n"
            description += df.head(100).to_string()
            chunks.append(DocumentChunk(
                content=description,
                source_path=filename,
                source_type=DocumentType.TABLE,
                chunk_index=0,
                metadata={"rows": len(df), "columns": len(df.columns), "columns_list": list(df.columns.astype(str))},
            ))
            return chunks
        except ImportError:
            return await self._fallback_text(file_path, filename, DocumentType.TABLE)

    async def _parse_code(self, file_path: str, filename: str) -> list[DocumentChunk]:
        import asyncio, json

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            source_code = f.read()

        language = detect_language(filename)
        total_lines = len(source_code.splitlines())
        chunks = []

        # ==========================================================
        # 通道一：静态分析——精确的事实提取
        # ==========================================================
        logger.info(f"[代码增强] 静态分析: {filename} ({language})")
        static_result = await asyncio.to_thread(analyze_code_file, file_path)

        functions = static_result.get("functions", [])
        classes = static_result.get("classes", [])
        imports = static_result.get("imports", [])
        complexity = static_result.get("complexity", [])
        call_graph = static_result.get("call_graph", {})

        # ==========================================================
        # 通道二：大模型语义总结——理解代码意图
        # ==========================================================
        semantic = {}
        if len(source_code) > 20:
            try:
                from openai import AsyncOpenAI
                client = AsyncOpenAI(
                    api_key=self.settings.deepseek_api_key,
                    base_url=self.settings.deepseek_base_url,
                )
                prompt = CODE_SEMANTIC_PROMPT.format(code=source_code[:8000])
                response = await client.chat.completions.create(
                    model=self.settings.deepseek_model,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.2,
                    max_tokens=2048,
                )
                resp_text = response.choices[0].message.content or ""
                match = re.search(r'\{.*\}', resp_text, re.DOTALL)
                if match:
                    semantic = json.loads(match.group())
                logger.info(f"[代码增强] 语义理解完成: {filename}, 摘要={semantic.get('summary','?')[:60]}")
            except Exception as e:
                logger.warning(f"[代码增强] DeepSeek 调用失败: {e}")

        # ==========================================================
        # 构建结构化 Chunk：架构概览
        # ==========================================================
        overview_parts = []
        overview_parts.append(f"=== 代码静态分析报告 ===")
        overview_parts.append(f"文件: {filename}")
        overview_parts.append(f"语言: {language} | 行数: {total_lines} | 函数: {len(functions)} | 类: {len(classes)} | 导入: {len(imports)}")

        if semantic.get("summary"):
            overview_parts.append(f"\n[功能摘要] {semantic['summary']}")
        if semantic.get("algorithm"):
            overview_parts.append(f"[核心算法] {semantic['algorithm']}")
        if semantic.get("business_logic"):
            overview_parts.append("\n[业务逻辑]")
            for i, step in enumerate(semantic["business_logic"], 1):
                overview_parts.append(f"  {i}. {step}")
        if semantic.get("inputs"):
            overview_parts.append(f"\n[输入] {semantic['inputs']}")
        if semantic.get("outputs"):
            overview_parts.append(f"[输出] {semantic['outputs']}")
        if semantic.get("edge_cases"):
            overview_parts.append(f"\n[边界情况]")
            for ec in semantic["edge_cases"]:
                overview_parts.append(f"  - {ec}")
        if semantic.get("related_concepts"):
            overview_parts.append(f"\n[相关概念] {', '.join(semantic['related_concepts'])}")
        if semantic.get("code_quality_notes"):
            overview_parts.append(f"\n[质量评价] {semantic['code_quality_notes']}")

        # 导入依赖
        if imports:
            overview_parts.append("\n[依赖导入]")
            for imp in imports[:30]:
                overview_parts.append(f"  - {json.dumps(imp, ensure_ascii=False)}")

        # 调用关系
        if call_graph.get("edges"):
            overview_parts.append("\n[内部调用关系]")
            for e in call_graph["edges"][:20]:
                overview_parts.append(f"  {e['from']} -> {e['to']}")

        chunks.append(DocumentChunk(
            content="\n".join(overview_parts),
            source_path=filename,
            source_type=DocumentType.CODE,
            chunk_index=0,
            metadata={
                "language": language,
                "lines": total_lines,
                "function_count": len(functions),
                "class_count": len(classes),
                "import_count": len(imports),
                "algorithm": semantic.get("algorithm"),
                "concepts": semantic.get("related_concepts", []),
                "extraction_mode": "enhanced_code",
                "chunk_type": "overview",
            },
        ))

        # ==========================================================
        # 构建结构化 Chunk：逐函数理解
        # ==========================================================
        for fi, func in enumerate(functions):
            func_parts = []
            func_parts.append(f"=== 函数: {func['name']} (第{func.get('lineno','?')}行) ===")

            if func.get("docstring"):
                func_parts.append(f"文档: {func['docstring'][:300]}")
            if "args" in func:
                func_parts.append(f"参数: {', '.join(func['args']) if func['args'] else '(无)'} (共{func.get('arg_count',0)}个)")
            if func.get("return_type"):
                func_parts.append(f"返回类型: {func['return_type']}")
            if func.get("decorators"):
                func_parts.append(f"装饰器: {', '.join(func['decorators'])}")
            if func.get("called_functions"):
                func_parts.append(f"内部调用: {', '.join(func['called_functions'][:15])}")
            if func.get("line_count"):
                func_parts.append(f"代码行数: {func['line_count']}")

            features = []
            if func.get("has_loop"):
                features.append("含循环")
            if func.get("has_condition"):
                features.append("含分支")
            if func.get("has_try"):
                features.append("含异常处理")
            if func.get("async"):
                features.append("异步")
            if features:
                func_parts.append(f"特征: {', '.join(features)}")

            # 复杂度
            for c in complexity:
                if c["name"] == func["name"]:
                    func_parts.append(f"圈复杂度: {c['complexity']} (等级: {c['rank']})")

            chunks.append(DocumentChunk(
                content="\n".join(func_parts),
                source_path=filename,
                source_type=DocumentType.CODE,
                chunk_index=fi + 1,
                metadata={
                    "language": language,
                    "function_name": func["name"],
                    "lineno": func.get("lineno", 0),
                    "line_count": func.get("line_count", 0),
                    "arg_count": func.get("arg_count", 0),
                    "has_loop": func.get("has_loop", False),
                    "has_condition": func.get("has_condition", False),
                    "extraction_mode": "enhanced_code",
                    "chunk_type": "function",
                },
            ))

        # ==========================================================
        # 构建结构化 Chunk：逐类理解
        # ==========================================================
        for ci, cls in enumerate(classes):
            cls_parts = []
            cls_parts.append(f"=== 类: {cls['name']} (第{cls.get('lineno','?')}行) ===")
            if cls.get("docstring"):
                cls_parts.append(f"文档: {cls['docstring'][:300]}")
            if cls.get("bases"):
                cls_parts.append(f"继承: {', '.join(cls['bases'])}")
            if cls.get("decorators"):
                cls_parts.append(f"装饰器: {', '.join(cls['decorators'])}")
            methods = cls.get("methods", [])
            if methods:
                cls_parts.append(f"方法 ({len(methods)}): {', '.join(methods)}")
            if cls.get("line_count"):
                cls_parts.append(f"代码行数: {cls['line_count']}")

            chunks.append(DocumentChunk(
                content="\n".join(cls_parts),
                source_path=filename,
                source_type=DocumentType.CODE,
                chunk_index=len(functions) + ci + 1,
                metadata={
                    "language": language,
                    "class_name": cls["name"],
                    "method_count": len(methods),
                    "lineno": cls.get("lineno", 0),
                    "extraction_mode": "enhanced_code",
                    "chunk_type": "class",
                },
            ))

        # 如果无结构化信息，保留原始代码文本作为 fallback
        if len(chunks) == 0:
            chunks.append(DocumentChunk(
                content=source_code[:5000],
                source_path=filename,
                source_type=DocumentType.CODE,
                chunk_index=0,
                metadata={"language": language, "lines": total_lines, "extraction_mode": "fallback"},
            ))

        logger.info(
            f"[代码增强] 完成: {filename}, {len(chunks)} chunks "
            f"(概览+{len(functions)}函数+{len(classes)}类), 语义={bool(semantic)}"
        )
        return chunks

    async def _parse_code_with_analysis(self, file_path: str, filename: str) -> tuple[list[DocumentChunk], dict]:
        import asyncio

        language = detect_language(filename)

        logger.info(f"[代码增强] 静态分析: {filename} ({language})")
        static_result = await asyncio.to_thread(analyze_code_file, file_path)

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            source_code = f.read()

        semantic = {}
        if len(source_code) > 20:
            try:
                from openai import AsyncOpenAI
                client = AsyncOpenAI(
                    api_key=self.settings.deepseek_api_key,
                    base_url=self.settings.deepseek_base_url,
                )
                prompt = CODE_SEMANTIC_PROMPT.format(code=source_code[:8000])
                response = await client.chat.completions.create(
                    model=self.settings.deepseek_model,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.2,
                    max_tokens=2048,
                )
                resp_text = response.choices[0].message.content or ""
                match = re.search(r'\{.*\}', resp_text, re.DOTALL)
                if match:
                    semantic = json.loads(match.group())
            except Exception as e:
                logger.warning(f"[代码增强] DeepSeek 调用失败: {e}")

        chunks = self._build_code_chunks(file_path, filename, source_code, language, static_result, semantic)

        code_analysis = {
            "file_path": file_path,
            "filename": filename,
            "language": language,
            "static_result": static_result,
            "semantic": {
                "summary": semantic.get("summary", ""),
                "algorithm": semantic.get("algorithm"),
                "business_logic": semantic.get("business_logic", []),
                "inputs": semantic.get("inputs", ""),
                "outputs": semantic.get("outputs", ""),
                "edge_cases": semantic.get("edge_cases", []),
                "related_concepts": semantic.get("related_concepts", []),
                "code_quality_notes": semantic.get("code_quality_notes", ""),
            },
        }

        return chunks, code_analysis

    def _build_code_chunks(
        self, file_path: str, filename: str, source_code: str,
        language: str, static_result: dict, semantic: dict,
    ) -> list[DocumentChunk]:
        import json

        total_lines = len(source_code.splitlines())
        functions = static_result.get("functions", [])
        classes = static_result.get("classes", [])
        imports = static_result.get("imports", [])
        complexity = static_result.get("complexity", [])
        call_graph = static_result.get("call_graph", {})
        chunks = []

        overview_parts = []
        overview_parts.append(f"=== 代码静态分析报告 ===")
        overview_parts.append(f"文件: {filename}")
        overview_parts.append(f"语言: {language} | 行数: {total_lines} | 函数: {len(functions)} | 类: {len(classes)} | 导入: {len(imports)}")

        if semantic.get("summary"):
            overview_parts.append(f"\n[功能摘要] {semantic['summary']}")
        if semantic.get("algorithm"):
            overview_parts.append(f"[核心算法] {semantic['algorithm']}")
        if semantic.get("business_logic"):
            overview_parts.append("\n[业务逻辑]")
            for i, step in enumerate(semantic["business_logic"], 1):
                overview_parts.append(f"  {i}. {step}")
        if semantic.get("inputs"):
            overview_parts.append(f"\n[输入] {semantic['inputs']}")
        if semantic.get("outputs"):
            overview_parts.append(f"[输出] {semantic['outputs']}")
        if semantic.get("edge_cases"):
            overview_parts.append(f"\n[边界情况]")
            for ec in semantic["edge_cases"]:
                overview_parts.append(f"  - {ec}")
        if semantic.get("related_concepts"):
            overview_parts.append(f"\n[相关概念] {', '.join(semantic['related_concepts'])}")
        if semantic.get("code_quality_notes"):
            overview_parts.append(f"\n[质量评价] {semantic['code_quality_notes']}")

        if imports:
            overview_parts.append("\n[依赖导入]")
            for imp in imports[:30]:
                overview_parts.append(f"  - {json.dumps(imp, ensure_ascii=False)}")

        if call_graph.get("edges"):
            overview_parts.append("\n[内部调用关系]")
            for e in call_graph["edges"][:20]:
                overview_parts.append(f"  {e['from']} -> {e['to']}")

        chunks.append(DocumentChunk(
            content="\n".join(overview_parts),
            source_path=filename,
            source_type=DocumentType.CODE,
            chunk_index=0,
            metadata={
                "language": language,
                "lines": total_lines,
                "function_count": len(functions),
                "class_count": len(classes),
                "import_count": len(imports),
                "algorithm": semantic.get("algorithm"),
                "concepts": semantic.get("related_concepts", []),
                "extraction_mode": "enhanced_code",
                "chunk_type": "overview",
            },
        ))

        for fi, func in enumerate(functions):
            func_parts = []
            func_parts.append(f"=== 函数: {func['name']} (第{func.get('lineno','?')}行) ===")
            if func.get("docstring"):
                func_parts.append(f"文档: {func['docstring'][:300]}")
            if "args" in func:
                func_parts.append(f"参数: {', '.join(func['args']) if func['args'] else '(无)'} (共{func.get('arg_count',0)}个)")
            if func.get("return_type"):
                func_parts.append(f"返回类型: {func['return_type']}")
            if func.get("decorators"):
                func_parts.append(f"装饰器: {', '.join(func['decorators'])}")
            if func.get("called_functions"):
                func_parts.append(f"内部调用: {', '.join(func['called_functions'][:15])}")
            if func.get("line_count"):
                func_parts.append(f"代码行数: {func['line_count']}")

            features = []
            if func.get("has_loop"):
                features.append("含循环")
            if func.get("has_condition"):
                features.append("含分支")
            if func.get("has_try"):
                features.append("含异常处理")
            if func.get("async"):
                features.append("异步")
            if features:
                func_parts.append(f"特征: {', '.join(features)}")

            for c in complexity:
                if c["name"] == func["name"]:
                    func_parts.append(f"圈复杂度: {c['complexity']} (等级: {c['rank']})")

            chunks.append(DocumentChunk(
                content="\n".join(func_parts),
                source_path=filename,
                source_type=DocumentType.CODE,
                chunk_index=fi + 1,
                metadata={
                    "language": language,
                    "function_name": func["name"],
                    "lineno": func.get("lineno", 0),
                    "line_count": func.get("line_count", 0),
                    "arg_count": func.get("arg_count", 0),
                    "has_loop": func.get("has_loop", False),
                    "has_condition": func.get("has_condition", False),
                    "extraction_mode": "enhanced_code",
                    "chunk_type": "function",
                },
            ))

        for ci, cls in enumerate(classes):
            cls_parts = []
            cls_parts.append(f"=== 类: {cls['name']} (第{cls.get('lineno','?')}行) ===")
            if cls.get("docstring"):
                cls_parts.append(f"文档: {cls['docstring'][:300]}")
            if cls.get("bases"):
                cls_parts.append(f"继承: {', '.join(cls['bases'])}")
            if cls.get("decorators"):
                cls_parts.append(f"装饰器: {', '.join(cls['decorators'])}")
            methods = cls.get("methods", [])
            if methods:
                cls_parts.append(f"方法 ({len(methods)}): {', '.join(methods)}")
            if cls.get("line_count"):
                cls_parts.append(f"代码行数: {cls['line_count']}")

            chunks.append(DocumentChunk(
                content="\n".join(cls_parts),
                source_path=filename,
                source_type=DocumentType.CODE,
                chunk_index=len(functions) + ci + 1,
                metadata={
                    "language": language,
                    "class_name": cls["name"],
                    "method_count": len(methods),
                    "lineno": cls.get("lineno", 0),
                    "extraction_mode": "enhanced_code",
                    "chunk_type": "class",
                },
            ))

        if len(chunks) == 0:
            chunks.append(DocumentChunk(
                content=source_code[:5000],
                source_path=filename,
                source_type=DocumentType.CODE,
                chunk_index=0,
                metadata={"language": language, "lines": total_lines, "extraction_mode": "fallback"},
            ))

        logger.info(
            f"[代码增强] 完成: {filename}, {len(chunks)} chunks "
            f"(概览+{len(functions)}函数+{len(classes)}类), 语义={bool(semantic)}"
        )
        return chunks

    async def _parse_text(self, file_path: str, filename: str) -> list[DocumentChunk]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        chunks = [DocumentChunk(
            content=content,
            source_path=filename,
            source_type=DocumentType.TEXT,
            chunk_index=0,
        )]
        return chunks

    async def _ensure_ffmpeg_path(self):
        import os, shutil
        try:
            import imageio_ffmpeg
            ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
            ffmpeg_dir = os.path.dirname(ffmpeg_exe)
            ffmpeg_standard = os.path.join(ffmpeg_dir, "ffmpeg.exe")
            if not os.path.exists(ffmpeg_standard):
                shutil.copy2(ffmpeg_exe, ffmpeg_standard)
            if ffmpeg_dir not in os.environ.get("PATH", ""):
                os.environ["PATH"] = ffmpeg_dir + os.pathsep + os.environ.get("PATH", "")
        except ImportError:
            pass

    async def _parse_video(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            await self._ensure_ffmpeg_path()
            import asyncio

            logger.info(f"[视频增强] 开始语音转写: {filename}")

            speech_result = await self.speech_processor.process_file(file_path, language="zh")
            text = speech_result.get("full_text", "").strip()

            if not text:
                logger.warning(f"视频转写结果为空: {filename}")
                raise ValueError("无语音内容")

            from zhconv import convert
            raw_segments = speech_result.get("segments", [])
            conf_report = speech_result.get("confidence_report", {})
            low_conf_words = speech_result.get("low_confidence_words", [])
            chunks = []

            # Step 1: 场景分割 + 关键帧提取
            logger.info(f"[视频增强] 场景检测: {filename}")
            scene_times = await asyncio.to_thread(
                self.enhancer.detect_scenes, file_path
            )
            logger.info(f"[视频增强] 检测到 {len(scene_times)} 个场景点: {filename}")

            temp_frames_dir = os.path.join(tempfile.gettempdir(), f"vid_frames_{uuid.uuid4().hex[:8]}")
            frame_paths = await asyncio.to_thread(
                self.enhancer.extract_keyframes, file_path, scene_times, temp_frames_dir
            )
            frame_descriptions = self.enhancer.describe_frames(frame_paths)
            try:
                for fp in frame_paths:
                    if os.path.exists(fp):
                        os.remove(fp)
                if os.path.exists(temp_frames_dir):
                    os.rmdir(temp_frames_dir)
            except Exception:
                pass

            # Step 2: 场景-语音事件绑定
            events = self.enhancer.bind_scene_events(raw_segments, scene_times, frame_descriptions)
            logger.info(f"[视频增强] 绑定 {len(events)} 个场景事件: {filename}")

            # Step 3: 结构化知识提取
            video_knowledge = await self.enhancer.extract_video_knowledge(events, filename)
            logger.info(f"[视频增强] 提取 {len(video_knowledge)} 条知识: {filename}")

            # Step 4: 构建增强chunks（按场景分段）
            scene_groups = {}
            for e in events:
                sidx = e.get("scene_index", 0)
                if sidx not in scene_groups:
                    scene_groups[sidx] = []
                scene_groups[sidx].append(e)

            chunk_idx = 0
            for sidx in sorted(scene_groups.keys()):
                group_events = scene_groups[sidx]
                first = group_events[0]
                last = group_events[-1]
                span_start = first["start"]
                span_end = last["end"]
                ts_range = f"[{int(span_start//60):02d}:{int(span_start%60):02d}-{int(span_end//60):02d}:{int(span_end%60):02d}]"

                event_texts = []
                for e in group_events:
                    ts = f"[{int(e['start']//60):02d}:{int(e['start']%60):02d}]"
                    txt = convert(e["text"].strip(), "zh-cn") if e["text"] else ""
                    event_texts.append(f"{ts} {txt}")

                frame_desc = first.get("frame_description", "")
                content_parts = [
                    f"=== 场景 {sidx + 1} {ts_range} ===",
                    "[语音转写]",
                    "\n".join(event_texts),
                ]
                if frame_desc:
                    content_parts.append(f"[画面描述]\n{frame_desc[:1000]}")

                # 附加上下文知识
                related_knowledge = [
                    k for k in video_knowledge
                    if k.get("evidence_ref", "") == ts_range
                ][:5]
                if related_knowledge:
                    k_lines = ["\n[场景相关知识]"]
                    for k_item in related_knowledge:
                        speaker = k_item.get("speaker", "")
                        fact = k_item.get("fact", "")
                        conf = k_item.get("confidence", 0.5)
                        k_lines.append(
                            f"- [说话人{speaker}] {fact} (置信度:{conf:.2f})"
                        )
                    content_parts.append("\n".join(k_lines))

                content = "\n\n".join(content_parts)
                content = convert(content, "zh-cn")

                chunks.append(DocumentChunk(
                    content=content,
                    source_path=filename,
                    source_type=DocumentType.VIDEO,
                    chunk_index=chunk_idx,
                    metadata={
                        "language": "zh",
                        "duration_seconds": raw_segments[-1].get("end", 0) if raw_segments else 0,
                        "scene_count": len(scene_groups),
                        "scene_index": sidx + 1,
                        "time_range": ts_range,
                        "event_count": len(group_events),
                        "frame_description": bool(frame_desc),
                        "knowledge_count": len(related_knowledge),
                        "extraction_mode": "enhanced_v3",
                        "whisper_model": "large-v3",
                        "transcription_confidence": conf_report.get("avg_confidence", 0),
                        "low_confidence_word_count": conf_report.get("low_count", 0),
                        "audio_preprocessed": speech_result.get("preprocessing", {}).get("applied", False),
                        "term_corrected": speech_result.get("correction", {}).get("term_corrected", False),
                        "llm_corrected": speech_result.get("correction", {}).get("llm_corrected", False),
                    },
                ))
                chunk_idx += 1

            logger.info(
                f"[视频增强] 完成: {filename}, {len(chunks)} 场景chunks, "
                f"{len(raw_segments)} 语音段, {len(scene_times)} 场景, {len(video_knowledge)} 知识点"
            )
            return chunks

        except ImportError as e:
            logger.warning(f"转写依赖缺失: {e}")
        except Exception as e:
            logger.warning(f"视频处理失败: {filename} - {e}")

        return [DocumentChunk(
            content=f"[未提取到语音内容] {filename}",
            source_path=filename,
            source_type=DocumentType.VIDEO,
            chunk_index=0,
            confidence=0.0,
            metadata={"extracted": False, "reason": "语音转写失败或无语音内容"},
        )]

    async def _parse_audio(self, file_path: str, filename: str) -> list[DocumentChunk]:
        try:
            await self._ensure_ffmpeg_path()
            import asyncio

            logger.info(f"[音频增强] 开始语音转写: {filename}")

            speech_result = await self.speech_processor.process_file(file_path, language="zh")
            text = speech_result.get("full_text", "").strip()

            if not text:
                logger.warning(f"音频转写结果为空: {filename}")
                raise ValueError("无语音内容")

            from zhconv import convert
            text = convert(text, "zh-cn")
            raw_segments = speech_result.get("segments", [])
            conf_report = speech_result.get("confidence_report", {})
            low_conf_words = speech_result.get("low_confidence_words", [])
            duration = raw_segments[-1].get("end", 0) if raw_segments else 0
            chunks = []

            # Step 1: 智能分段
            seg_groups = self.enhancer.smart_segment(raw_segments)
            logger.info(f"[音频增强] 智能分段: {len(seg_groups)} 个语义段落")

            # Step 2: 逐段生成结构化摘要
            for sg_idx, sg in enumerate(seg_groups):
                sg_text_full = " ".join(s.get("text", "") for s in sg)
                sg_text = convert(sg_text_full, "zh-cn")
                sg_start = sg[0].get("start", 0)
                sg_end = sg[-1].get("end", 0)
                ts_range = f"[{int(sg_start//60):02d}:{int(sg_start%60):02d}-{int(sg_end//60):02d}:{int(sg_end%60):02d}]"

                summary = await self.enhancer.generate_segment_summary(sg)
                key_sents = self.enhancer.extract_key_sentences(sg_text, top_k=3)
                emotion = self.enhancer.detect_emotion(sg_text)
                intent = self.enhancer.detect_intent(sg_text)

                content_parts = [
                    f"=== 段落 {sg_idx + 1} {ts_range} ===",
                    f"标题: {summary.get('title', '无标题')}",
                    f"情绪: {emotion} | 意图: {intent}",
                    "",
                    "[核心要点]",
                ]
                for pt in summary.get("key_points", []):
                    content_parts.append(f"  - {convert(str(pt), 'zh-cn')}")
                if summary.get("keywords"):
                    content_parts.append(f"关键词: {' '.join(summary['keywords'])}")

                content_parts.append("")
                content_parts.append("[关键句]")
                for ks in key_sents:
                    content_parts.append(f"  * {ks}")

                content_parts.append("")
                content_parts.append("[全文转录]")
                content_parts.append(sg_text[:2000])

                content = convert("\n".join(content_parts), "zh-cn")
                chunks.append(DocumentChunk(
                    content=content,
                    source_path=filename,
                    source_type=DocumentType.AUDIO,
                    chunk_index=sg_idx,
                    metadata={
                        "language": "zh",
                        "duration_seconds": duration,
                        "segment_index": sg_idx + 1,
                        "segment_count": len(seg_groups),
                        "time_range": ts_range,
                        "title": convert(str(summary.get("title", "")), "zh-cn"),
                        "emotion": emotion,
                        "intent": intent,
                        "keywords": summary.get("keywords", []),
                        "key_sentences": key_sents[:3],
                        "extraction_mode": "enhanced_v3",
                        "whisper_model": "large-v3",
                        "transcription_confidence": conf_report.get("avg_confidence", 0),
                        "low_confidence_word_count": conf_report.get("low_count", 0),
                        "audio_preprocessed": speech_result.get("preprocessing", {}).get("applied", False),
                        "term_corrected": speech_result.get("correction", {}).get("term_corrected", False),
                        "llm_corrected": speech_result.get("correction", {}).get("llm_corrected", False),
                        "low_confidence_words_low_conf_ratio": conf_report.get("low_ratio", 0),
                    },
                ))

            # Step 3: 全局 Q&A 提取
            qa_pairs = await self.enhancer.extract_qa_pairs(text)
            if qa_pairs:
                qa_parts = ["=== 问答对提取 ==="]
                for qa in qa_pairs:
                    q = convert(str(qa.get("question", "")), "zh-cn")
                    a = convert(str(qa.get("answer", "")), "zh-cn")
                    qa_parts.append(f"Q: {q}")
                    qa_parts.append(f"A: {a}")
                    qa_parts.append("")
                qa_content = "\n".join(qa_parts)
                chunks.append(DocumentChunk(
                    content=qa_content,
                    source_path=filename,
                    source_type=DocumentType.AUDIO,
                    chunk_index=len(seg_groups),
                    metadata={
                        "language": "zh",
                        "duration_seconds": duration,
                        "qa_pair_count": len(qa_pairs),
                        "extraction_mode": "enhanced_v3",
                        "whisper_model": "large-v3",
                        "transcription_confidence": conf_report.get("avg_confidence", 0),
                    },
                ))

            logger.info(
                f"[音频增强] 完成: {filename}, {len(chunks)} chunks, "
                f"{len(seg_groups)} 段落, {len(qa_pairs)} Q&A"
            )
            return chunks

        except ImportError as e:
            logger.warning(f"转写依赖缺失: {e}")
        except Exception as e:
            logger.warning(f"音频处理失败: {filename} - {e}")

        return [DocumentChunk(
            content=f"[未提取到语音内容] {filename}",
            source_path=filename,
            source_type=DocumentType.AUDIO,
            chunk_index=0,
            confidence=0.0,
            metadata={"extracted": False, "reason": "语音转写失败或无语音内容"},
        )]

    async def _fallback_text(self, file_path: str, filename: str, file_type: DocumentType) -> list[DocumentChunk]:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        except Exception:
            content = f"[无法解析的文件] {filename}"
        return [DocumentChunk(
            content=content,
            source_path=filename,
            source_type=file_type,
            chunk_index=0,
            confidence=0.3,
        )]